"""Wrap the rendered slides into a PPTX, one full-bleed image per slide.

Called by build.mjs with the PNG paths in order.

Image slides rather than text boxes is deliberate: it means the deck looks
identical on any machine, because there are no fonts to substitute. That
is the usual way a deck breaks minutes before a pitch. The cost is that
the text is not editable in PowerPoint — edit the HTML in build.mjs and
rebuild.
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

# 1920x1080 at 96 DPI in English Metric Units — a true 16:9 slide.
SLIDE_W = Emu(1920 * 9525)
SLIDE_H = Emu(1080 * 9525)


def main(pngs: list[str]) -> None:
    if not pngs:
        raise SystemExit("no slides given")

    deck = Presentation()
    deck.slide_width = SLIDE_W
    deck.slide_height = SLIDE_H
    blank = deck.slide_layouts[6]  # blank: no placeholders to fight with

    for png in pngs:
        slide = deck.slides.add_slide(blank)
        slide.shapes.add_picture(png, 0, 0, width=SLIDE_W, height=SLIDE_H)

    out = Path(pngs[0]).parent.parent / "divergence-deck.pptx"
    deck.save(out)

    size_mb = out.stat().st_size / 1e6
    print(f"{out}")
    print(f"  {len(pngs)} slides · 1920x1080 · {size_mb:.1f} MB")


if __name__ == "__main__":
    main(sys.argv[1:])
